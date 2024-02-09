import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor  # Corrected import statement

# Load data from Excel (adjust the file path accordingly)
data = pd.read_excel(r'C:\Users\AnkitKumar\Desktop\data.xlsx')

# Create a PowerPoint presentation with 16:9 aspect ratio
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Function to add a slide with specified layout
def add_slide(store_name, address, date_of_audit, time_of_audit,store_address):
    # Use 16:9 blank slide layout
    slide_layout = presentation.slide_layouts[5]  

    # Add a new slide
    slide = presentation.slides.add_slide(slide_layout)

    # Add title with store name
    title = slide.shapes.title
    title.text = f"INSTANT-INK"

    # Set font for the title
    title.text_frame.text = f"INSTANT INK - INSTORE CX"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Blue color
    title.text_frame.paragraphs[0].font.name = "Forma DJR Display (Headings)"  # Set the font name

    title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # Add a table covering 20% of the slide from the top
    left = Inches(0.65)
    top = Inches(1.2)
    width = Inches(14.4)
    height = Inches(0.9)

    table = slide.shapes.add_table(rows=3, cols=2, left=left, top=top, width=width, height=height).table
    table.cell(0, 0).text = 'Retailer'
    table.cell(0, 1).text = store_name
    table.cell(1, 0).text = 'Date & Time of Audit:'
    table.cell(1, 1).text = f"{date_of_audit} - {time_of_audit}"
    table.cell(2, 0).text = 'Store Address:'
    table.cell(2, 1).text = store_address
    

    # Set font and color for the table
    for row in table.rows:
        for cell in row.cells:
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black color

     

     # Set background color for the cells in the first row to black
    for cell in table.rows[0].cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black color
            

     # Set font color of the text in the first row to black
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(255, 255, 255)  # Black color



  #insert picture 

    # logo_path = r'C:\Users\AnkitKumar\Desktop\logo.jpg'  # Provide the correct file path for your logo
    # left_logo = Inches(14)  # Adjust this value to position the logo more to the right
    # top_logo = Inches(9)   # Adjust this value to position the logo towards the bottom
    # logo_width = Inches(2)
    # logo_height = Inches(0.5)
    # slide.shapes.add_picture(logo_path, left_logo, top_logo, width=logo_width, height=logo_height)



# Iterate through each row in the data table
for index, row in data.iterrows():
    print("inside loop")
    add_slide(row['Store Name'], row['Address'], row['Date of Audit'], row['Time of Audit'],row['Store Address'])

# Save the presentation
presentation.save(r'C:\Users\AnkitKumar\Desktop\output_presentation.pptx')
